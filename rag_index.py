"""
Persistent FAISS index management for the AWS Lambda RAG statement rephraser.

Why this module exists
----------------------
Building a FAISS index requires downloading the source document (PDF or PPTX)
from S3, parsing it, splitting it into chunks, and embedding every chunk with
Bedrock. Doing that on *every* Lambda invocation is slow and expensive.

Instead we split the work into two phases:

1. BUILD (run on a schedule): download the source document -> split -> embed ->
   build FAISS index -> serialize the index to S3. This is done by the
   standalone `build_embedding.py` script (see `build_index()` there).

2. QUERY (every invocation): load the serialized FAISS index from S3 once per
   warm container, then run similarity search + LLM generation to REPHRASE the
   incoming statement following the style/format of the Job Aids document. See
   `load_index()` and `query_index()`.

The FAISS index is stored in S3 as a single pickle file. Because the index is
read-only after build, it is safe to load it once and cache it in a module-level
global so warm Lambda invocations skip the S3 download entirely.
"""
import io
import json
import os
import pickle
import tempfile

import boto3
from langchain_community.document_loaders import PyPDFLoader
from langchain_community.embeddings import BedrockEmbeddings
from langchain_community.llms import Bedrock
from langchain_community.vectorstores import FAISS
from langchain_text_splitters import RecursiveCharacterTextSplitter

# ---------------------------------------------------------------------------
# Configuration (all overridable via Lambda environment variables)
# ---------------------------------------------------------------------------
S3_BUCKET = os.environ.get("S3_BUCKET", "dealerplatform")
# Default source document is the 2025 Dealer Profile Job Aid (PowerPoint).
SOURCE_KEY = os.environ.get("S3_OBJECT_KEY", "2025 Dealer Profile Job Aid.pptx")
INDEX_KEY = os.environ.get("S3_INDEX_KEY", "rag/faiss_index.pkl")

EMBEDDING_MODEL = os.environ.get("EMBEDDING_MODEL", "amazon.titan-embed-text-v1")
# NOTE: get_llm() uses the Amazon Nova Messages API request/response format
# (inferenceConfig + output.message.content). The LLM_MODEL MUST therefore be an
# Amazon Nova model (amazon.nova-*). A Claude model id would receive a Nova-style
# body and fail. Override via the LLM_MODEL env var if you prefer a different
# Nova model id (e.g. amazon.nova-pro-v1:0).
LLM_MODEL = os.environ.get("LLM_MODEL", "amazon.nova-lite-v1:0")

CHUNK_SIZE = int(os.environ.get("CHUNK_SIZE", "1000"))
CHUNK_OVERLAP = int(os.environ.get("CHUNK_OVERLAP", "200"))

# Module-level cache so warm Lambda containers reuse the loaded index.
_INDEX_CACHE = None


# ---------------------------------------------------------------------------
# Document loader (supports both PDF and PPTX)
# ---------------------------------------------------------------------------
def load_document(path):
    """Load a local document into LangChain Documents.

    The loader is chosen based on the file extension:
      - .pdf  -> PyPDFLoader
      - .pptx / .ppt -> lightweight python-pptx based loader (no `unstructured`
        dependency, keeping the Lambda layer small)

    Returns a list of LangChain Document objects.
    """
    ext = os.path.splitext(path)[1].lower()

    if ext == ".pdf":
        loader = PyPDFLoader(path)
        return loader.load()

    if ext in (".pptx", ".ppt"):
        return _load_pptx(path)

    raise ValueError(
        f"Unsupported document type '{ext}'. Supported types: .pdf, .pptx, .ppt"
    )


def _load_pptx(path):
    """Extract text from a PowerPoint file using python-pptx directly.

    Returns a list of LangChain Document objects, one per slide. This avoids
    the heavy `unstructured` dependency so the Lambda layer stays small.
    """
    from langchain_core.documents import Document
    from pptx import Presentation

    prs = Presentation(path)
    docs = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    row_text = " | ".join(c for c in cells if c)
                    if row_text:
                        texts.append(row_text)
        content = "\n".join(texts).strip()
        if content:
            docs.append(
                Document(page_content=content, metadata={"slide": idx})
            )
    return docs


# ---------------------------------------------------------------------------
# Embeddings / LLM helpers
# ---------------------------------------------------------------------------
def get_embeddings():
    """Return a BedrockEmbeddings client. Credentials come from the Lambda
    execution role (or the default AWS profile when run locally)."""
    return BedrockEmbeddings(model_id=EMBEDDING_MODEL)


def get_llm():
    """Return a callable that generates an answer using the Bedrock LLM.

    Uses the Amazon Nova Messages API directly via boto3 so it works with
    Amazon-owned models (which do not require the Anthropic use-case approval)
    and with the current inference-profile model IDs.
    """
    client = boto3.client("bedrock-runtime")

    def _invoke(prompt):
        body = {
            "inferenceConfig": {
                "max_new_tokens": 300,
                "temperature": 0.1,
                "top_p": 0.9,
            },
            "messages": [
                {"role": "user", "content": [{"text": prompt}]},
            ],
        }
        response = client.invoke_model(
            modelId=LLM_MODEL,
            contentType="application/json",
            accept="application/json",
            body=json.dumps(body),
        )
        payload = json.loads(response["body"].read())
        # Nova returns the answer under output.message.content[0].text
        return payload["output"]["message"]["content"][0]["text"]

    return _invoke


# ---------------------------------------------------------------------------
# BUILD phase (run on a schedule via build_embedding.py)
# ---------------------------------------------------------------------------
def build_index(bucket=None, source_key=None, index_key=None):
    """Download the source document (PDF or PPTX) from S3, split + embed it,
    build a FAISS index, and serialize it back to S3.

    This is normally invoked by the standalone `build_embedding.py` script on a
    schedule so the index is rebuilt whenever the source document changes.

    Returns the FAISS vectorstore (also cached in memory).
    """
    bucket = bucket or S3_BUCKET
    source_key = source_key or SOURCE_KEY
    index_key = index_key or INDEX_KEY

    s3 = boto3.client("s3")

    # 1. Download the source document to a local temp file.
    #    The loaders need a local path with the correct extension so they can
    #    detect the file type.
    ext = os.path.splitext(source_key)[1] or ".pdf"
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        s3.download_fileobj(bucket, source_key, tmp)
        tmp_path = tmp.name

    try:
        # 2. Load + split the document.
        documents = load_document(tmp_path)

        splitter = RecursiveCharacterTextSplitter(
            separators=["\n\n", "\n", " ", ""],
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP,
        )
        chunks = splitter.split_documents(documents)

        if not chunks:
            raise ValueError(
                f"No text content extracted from s3://{bucket}/{source_key}. "
                "The document may be empty or image-only."
            )

        # 3. Embed + build the FAISS index.
        vectorstore = FAISS.from_documents(chunks, get_embeddings())

        # 4. Serialize the index and upload it to S3.
        # The FAISS vectorstore retains a reference to the embedding function,
        # which wraps a non-picklable boto3 Bedrock client. Strip it before
        # pickling and restore it afterwards so the index can be serialized.
        embedding_fn = vectorstore.embedding_function
        vectorstore.embedding_function = None
        try:
            buffer = io.BytesIO()
            pickle.dump(vectorstore, buffer)
            buffer.seek(0)
            s3.upload_fileobj(buffer, bucket, index_key)
        finally:
            vectorstore.embedding_function = embedding_fn

        # Cache it so the same process can query immediately.
        global _INDEX_CACHE
        _INDEX_CACHE = vectorstore

        print(f"Index built and uploaded to s3://{bucket}/{index_key} "
              f"({len(chunks)} chunks from {source_key}).")
        return vectorstore
    finally:
        # Clean up the temp file.
        try:
            os.remove(tmp_path)
        except OSError:
            pass


# ---------------------------------------------------------------------------
# QUERY phase (every invocation)
# ---------------------------------------------------------------------------
def load_index(bucket=None, index_key=None, force_reload=False):
    """Load the FAISS index from S3, caching it for warm invocations.

    Returns the FAISS vectorstore.
    """
    global _INDEX_CACHE
    if _INDEX_CACHE is not None and not force_reload:
        return _INDEX_CACHE

    bucket = bucket or S3_BUCKET
    index_key = index_key or INDEX_KEY

    s3 = boto3.client("s3")
    buffer = io.BytesIO()
    s3.download_fileobj(bucket, index_key, buffer)
    buffer.seek(0)

    _INDEX_CACHE = pickle.load(buffer)
    # Restore the embedding function that was stripped before pickling.
    if getattr(_INDEX_CACHE, "embedding_function", None) is None:
        _INDEX_CACHE.embedding_function = get_embeddings()
    print(f"Index loaded from s3://{bucket}/{index_key}.")
    return _INDEX_CACHE


def query_index(statement, vectorstore=None, k=4):
    """Rephrase an incoming statement following the style/format of the Job Aids
    document.

    The Job Aids index is loaded from S3 (cached for warm invocations), the most
    relevant chunks are retrieved via similarity search, and the Bedrock LLM is
    prompted to rephrase the statement the way the Job Aids document suggests.

    Returns a dict with the rephrased statement and the retrieved source chunks.
    """
    vectorstore = vectorstore or load_index()

    # 1. Similarity search: find the k most relevant Job Aids chunks.
    docs = vectorstore.similarity_search(statement, k=k)

    # 2. Build a rephrasing prompt from the retrieved Job Aids chunks.
    context = "\n\n".join(doc.page_content for doc in docs)
    prompt = (
        "You are an expert at rephrasing statements following the style and "
        "format of a Job Aids reference document. Read the reference content "
        "below, understand how it phrases and structures statements, then "
        "rephrase the given statement in that same style. Keep the meaning "
        "unchanged; only change the wording, tone, and structure to match the "
        "reference. Output ONLY the rephrased statement.\n\n"
        "REFERENCE (Job Aids content):\n"
        f"{context}\n\n"
        f"STATEMENT TO REPHRASE:\n{statement}\n\n"
        "REPHRASED STATEMENT:"
    )

    # 3. Generate the rephrased statement with the Bedrock LLM.
    llm = get_llm()
    rephrased = llm(prompt)

    return {
        "rephrased_statement": rephrased,
        "sources": [doc.page_content for doc in docs],
    }


# ---------------------------------------------------------------------------
# CLI entry point for the BUILD phase
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    # The index build is normally done by the standalone build_embedding.py
    # script on a schedule. This block is kept for convenience:
    #   python rag_index.py
    build_index()
