"""
Standalone embeddings/index build script for the AWS Lambda RAG rephraser.

Why this file exists
--------------------
Building a FAISS index requires downloading the source document (PDF or PPTX)
from S3, parsing it, splitting it into chunks, and embedding every chunk with
Bedrock. Doing that on *every* Lambda invocation is slow and expensive.

Instead, this script builds the index ONCE and uploads the serialized FAISS
index to S3. It is designed to be run on a SCHEDULE (e.g. an EventBridge rule
invoking a build Lambda, or a cron job) so that whenever the source document is
updated, the index is rebuilt automatically.

The QUERY phase (every invocation) then only loads the pre-built index from S3
and runs similarity search + LLM generation. See `rag_index.py`.

Usage
-----
Run locally (with AWS credentials configured):
    python build_embedding.py

Or invoke as a Lambda handler (for an EventBridge schedule):
    handler = build_embedding.lambda_handler

Both PDF (.pdf) and PowerPoint (.pptx / .ppt) source documents are supported.
The loader is chosen automatically based on the file extension.
"""
import io
import json
import os
import pickle
import tempfile

import boto3
from langchain_community.document_loaders import PyPDFLoader
from langchain_community.embeddings import BedrockEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_text_splitters import RecursiveCharacterTextSplitter

# ---------------------------------------------------------------------------
# Configuration (all overridable via environment variables)
# ---------------------------------------------------------------------------
S3_BUCKET = os.environ.get("S3_BUCKET", "dealerplatform")
# Default source document is the 2025 Dealer Profile Job Aid (PowerPoint).
SOURCE_KEY = os.environ.get("S3_OBJECT_KEY", "2025 Dealer Profile Job Aid.pptx")
INDEX_KEY = os.environ.get("S3_INDEX_KEY", "rag/faiss_index.pkl")

EMBEDDING_MODEL = os.environ.get("EMBEDDING_MODEL", "amazon.titan-embed-text-v1")

CHUNK_SIZE = int(os.environ.get("CHUNK_SIZE", "1000"))
CHUNK_OVERLAP = int(os.environ.get("CHUNK_OVERLAP", "200"))


# ---------------------------------------------------------------------------
# Document loader (supports both PDF and PPTX)
# ---------------------------------------------------------------------------
def load_document(path):
    """Load a local document into LangChain Documents.

    The loader is chosen based on the file extension:
      - .pdf  -> PyPDFLoader
      - .pptx / .ppt -> lightweight python-pptx based loader (no `unstructured`
        dependency, keeping the Lambda layer small)

    Returns a list of LangChain Document objects.
    """
    ext = os.path.splitext(path)[1].lower()

    if ext == ".pdf":
        loader = PyPDFLoader(path)
        return loader.load()

    if ext in (".pptx", ".ppt"):
        return _load_pptx(path)

    raise ValueError(
        f"Unsupported document type '{ext}'. Supported types: .pdf, .pptx, .ppt"
    )


def _load_pptx(path):
    """Extract text from a PowerPoint file using python-pptx directly.

    Returns a list of LangChain Document objects, one per slide. This avoids
    the heavy `unstructured` dependency so the Lambda layer stays small.
    """
    from langchain_core.documents import Document
    from pptx import Presentation

    prs = Presentation(path)
    docs = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    row_text = " | ".join(c for c in cells if c)
                    if row_text:
                        texts.append(row_text)
        content = "\n".join(texts).strip()
        if content:
            docs.append(
                Document(page_content=content, metadata={"slide": idx})
            )
    return docs


def get_embeddings():
    """Return a BedrockEmbeddings client.

    Credentials come from the Lambda execution role (or the default AWS profile
    when run locally).
    """
    return BedrockEmbeddings(model_id=EMBEDDING_MODEL)


# ---------------------------------------------------------------------------
# Build phase
# ---------------------------------------------------------------------------
def build_index(bucket=None, source_key=None, index_key=None):
    """Download the source document from S3, split + embed it, build a FAISS
    index, and serialize it back to S3.

    Returns the FAISS vectorstore.
    """
    bucket = bucket or S3_BUCKET
    source_key = source_key or SOURCE_KEY
    index_key = index_key or INDEX_KEY

    s3 = boto3.client("s3")

    # 1. Download the source document to a local temp file.
    #    The loaders need a local path with the correct extension so they can
    #    detect the file type.
    ext = os.path.splitext(source_key)[1] or ".pdf"
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        s3.download_fileobj(bucket, source_key, tmp)
        tmp_path = tmp.name

    try:
        # 2. Load + split the document.
        documents = load_document(tmp_path)

        splitter = RecursiveCharacterTextSplitter(
            separators=["\n\n", "\n", " ", ""],
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP,
        )
        chunks = splitter.split_documents(documents)

        if not chunks:
            raise ValueError(
                f"No text content extracted from s3://{bucket}/{source_key}. "
                "The document may be empty or image-only."
            )

        # 3. Embed + build the FAISS index.
        vectorstore = FAISS.from_documents(chunks, get_embeddings())

        # 4. Serialize the index and upload it to S3.
        # The FAISS vectorstore retains a reference to the embedding function,
        # which wraps a non-picklable boto3 Bedrock client. Strip it before
        # pickling and restore it afterwards so the index can be serialized.
        embedding_fn = vectorstore.embedding_function
        vectorstore.embedding_function = None
        try:
            buffer = io.BytesIO()
            pickle.dump(vectorstore, buffer)
            buffer.seek(0)
            s3.upload_fileobj(buffer, bucket, index_key)
        finally:
            vectorstore.embedding_function = embedding_fn

        print(
            f"Index built and uploaded to s3://{bucket}/{index_key} "
            f"({len(chunks)} chunks from {source_key})."
        )
        return vectorstore
    finally:
        # Clean up the temp file.
        try:
            os.remove(tmp_path)
        except OSError:
            pass


# ---------------------------------------------------------------------------
# Lambda handler (for EventBridge schedule / scheduler-triggered runs)
# ---------------------------------------------------------------------------
def lambda_handler(event, context):
    """Lambda entry point so this script can be triggered on a schedule.

    Expected event (EventBridge schedule rule):
        {}   (empty event is fine - config comes from env vars)

    Or a plain JSON event that can override the config:
        {
            "S3_BUCKET": "dealerplatform",
            "S3_OBJECT_KEY": "2025 Dealer Profile Job Aid.pptx",
            "S3_INDEX_KEY": "rag/faiss_index.pkl"
        }
    """
    try:
        # Allow the event to override environment-variable config.
        if isinstance(event, dict):
            bucket = event.get("S3_BUCKET", S3_BUCKET)
            source_key = event.get("S3_OBJECT_KEY", SOURCE_KEY)
            index_key = event.get("S3_INDEX_KEY", INDEX_KEY)
        else:
            bucket, source_key, index_key = S3_BUCKET, SOURCE_KEY, INDEX_KEY

        build_index(bucket=bucket, source_key=source_key, index_key=index_key)

        return {
            "statusCode": 200,
            "body": json.dumps(
                {
                    "message": "Index rebuilt successfully.",
                    "source": source_key,
                    "index": f"s3://{bucket}/{index_key}",
                }
            ),
        }
    except Exception as exc:  # noqa: BLE001 - surface any error to the caller
        print(f"Index build failed: {exc}")
        return {
            "statusCode": 500,
            "body": json.dumps({"error": str(exc)}),
        }


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    # Run locally (or in a one-off Lambda) to build + persist the index:
    #   python build_embedding.py
    build_index()
